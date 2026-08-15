import os
import sys
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def generate_sih_presentation():
    input_path = r'D:\NEW CODING PROJECTS\SIH_INTERNAL\documents\SIH PPT FORMAT.pptx'
    if not os.path.exists(input_path):
        input_path = r'D:\NEW CODING PROJECTS\SIH_INTERNAL\SIH PPT FORMAT.pptx'

    output_pptx = r'D:\NEW CODING PROJECTS\SIH_INTERNAL\SIH_2026_AgriVision_AI_Submission.pptx'
    output_pdf = r'D:\NEW CODING PROJECTS\SIH_INTERNAL\SIH_2026_AgriVision_AI_Submission.pdf'
    
    doc_pptx = r'D:\NEW CODING PROJECTS\SIH_INTERNAL\documents\SIH_2026_AgriVision_AI_Submission.pptx'
    doc_pdf = r'D:\NEW CODING PROJECTS\SIH_INTERNAL\documents\SIH_2026_AgriVision_AI_Submission.pdf'

    prs = pptx.Presentation(input_path)
    
    # Helper to style top-left oval
    def style_oval(oval_shape):
        oval_shape.left = Inches(0.3)
        oval_shape.top = Inches(0.2)
        oval_shape.width = Inches(1.8)
        oval_shape.height = Inches(0.9)
        tf = oval_shape.text_frame
        tf.word_wrap = False
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "NeoMedtech"
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(16, 185, 129)

    # ----------------------------------------------------
    # SLIDE 1: TITLE PAGE
    # ----------------------------------------------------
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if 'SMART INDIA HACKATHON' in text:
                shape.text_frame.text = "SMART INDIA HACKATHON 2026"
                p = shape.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(28)
                p.font.color.rgb = RGBColor(16, 185, 129)
            elif 'TITLE PAGE' in text:
                shape.text_frame.text = "IIT JODHPUR INTERNAL HACKATHON EVALUATION"
                p = shape.text_frame.paragraphs[0]
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = RGBColor(100, 116, 139)
            elif 'Problem Statement ID' in text:
                shape.top = Inches(1.8)
                shape.left = Inches(0.4)
                shape.width = Inches(8.5)
                shape.height = Inches(4.8)
                tf = shape.text_frame
                tf.clear()
                
                details = [
                    ("Problem Statement ID:", " Problem Statement 1 (Software PS)"),
                    ("Problem Statement Title:", " Unified AI Agri-Vision Platform for Crop Advisory and Livestock Management"),
                    ("Theme:", " Agriculture, FoodTech & Rural Development"),
                    ("PS Category:", " Software"),
                    ("Team Name:", " NeoMedtech"),
                    ("Institute & AISHE Code:", " Indian Institute of Technology Jodhpur (IIT Jodhpur) [U-0395]"),
                    ("Team Composition:", " Maheshiv Prajapat (Lead, AIIMS & IITJ), Dr. Eshitaa Panwar, Siddhant Shenvi, Mrunal Sonawale, Kishore Vijayakumar, Megha")
                ]
                
                for idx, (label, val) in enumerate(details):
                    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                    p.space_after = Pt(6)
                    r1 = p.add_run()
                    r1.text = label
                    r1.font.bold = True
                    r1.font.size = Pt(13)
                    r1.font.color.rgb = RGBColor(15, 23, 42)
                    
                    r2 = p.add_run()
                    r2.text = val
                    r2.font.bold = False
                    r2.font.size = Pt(13)
                    r2.font.color.rgb = RGBColor(51, 65, 85)

    # ----------------------------------------------------
    # SLIDE 2: PROPOSED SOLUTION
    # ----------------------------------------------------
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_text_frame:
            if 'IDEA TITLE' in shape.text:
                shape.text_frame.text = "AgriVision AI: Unified Neural Vision & Clinical Veterinary Triage"
                p = shape.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(16, 185, 129)
            elif 'Your Team Name' in shape.text or 'NeoMe' in shape.text:
                style_oval(shape)
            elif 'Proposed Solution' in shape.text:
                shape.top = Inches(1.3)
                shape.left = Inches(0.6)
                shape.width = Inches(11.4)
                shape.height = Inches(5.2)
                tf = shape.text_frame
                tf.clear()
                
                bullets = [
                    ("1. Dual-Domain Diagnostic Engine:", " Directly solves PS-1 by combining foliar crop pathology with AIIMS Jodhpur-standard veterinary clinical triage on a single unified mobile platform."),
                    ("2. Dual-Tier Neural Vision Pipeline:", " Edge-optimized Multimodal Vision Transformer architecture with automated resilient fallback for instant disease detection & healthy leaf certification."),
                    ("3. Clinical Livestock Triage & 1962 Dispatch:", " Evidence-based triage protocols for Lumpy Skin Disease, Foot & Mouth Disease (FMD), and Mastitis with 1-tap 1962 Emergency Ambulance dialing."),
                    ("4. Precision Agro-Dosage Calculator:", " Computes exact tank water and chemical fungicide grams based on farmer's land size (Bigha/Acre/Hectare) and knapsack pump capacity."),
                    ("5. Multilingual Voice & Local Dialect AI:", " Native speech recognition and spoken audio readout in Hindi, Marwari, and English for low-literacy rural farmers."),
                    ("6. Historical Digital Farm Ledger:", " Integrated dairy milk yield logger, fertilizer expense records, and net seasonal profitability tracking for marginal farmers.")
                ]
                
                for idx, (title, desc) in enumerate(bullets):
                    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                    p.space_after = Pt(6)
                    r1 = p.add_run()
                    r1.text = title
                    r1.font.bold = True
                    r1.font.size = Pt(13)
                    r1.font.color.rgb = RGBColor(15, 23, 42)
                    
                    r2 = p.add_run()
                    r2.text = desc
                    r2.font.bold = False
                    r2.font.size = Pt(12)
                    r2.font.color.rgb = RGBColor(51, 65, 85)

    # ----------------------------------------------------
    # SLIDE 3: TECHNICAL APPROACH
    # ----------------------------------------------------
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.has_text_frame:
            if 'TECHNICAL APPROACH' in shape.text:
                shape.text_frame.text = "TECHNICAL APPROACH & SYSTEM ARCHITECTURE"
                p = shape.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(16, 185, 129)
            elif 'Your Team Name' in shape.text or 'NeoMe' in shape.text:
                style_oval(shape)
            elif 'Technologies to be used' in shape.text or '1. Frontend' in shape.text:
                shape.top = Inches(1.3)
                shape.left = Inches(0.6)
                shape.width = Inches(11.4)
                shape.height = Inches(5.2)
                tf = shape.text_frame
                tf.clear()
                
                tech_bullets = [
                    ("1. Frontend & Client UX:", " React.js (Vite), High-Contrast Sunlight Theme, Web Speech API (Voice Recognition & TTS), WebRTC Live Camera Stream, HTML5 Canvas."),
                    ("2. AI Vision & Multimodal Core:", " Dual-Tier Multimodal Neural Vision Architecture (Primary Vision Transformer + Automated Fallback Model) with structured clinical prompt schemas, severity scoring, and organic/chemical protocol synthesis."),
                    ("3. Live Meteorology & APMC Mandi Feeds:", " Real-time Open-Meteo Satellite API integration (live temp, humidity, wind & precipitation chance) + APMC e-NAM wholesale market ticker."),
                    ("4. Implementation Pipeline & Workflow:", " [Farmer Camera Snap / Upload] -> [Canvas Image Compression] -> [Dual-Tier Multimodal Inference] -> [Dosage & First-Aid Protocol Engine] -> [Voice Playback & PDF Export]."),
                    ("5. Production Cloud Deployment:", " Hosted on Vercel Global Edge Network with sub-1500ms latency, CI/CD automated via GitHub repository."),
                    ("6. Security & Offline Resilience:", " Zero frontend API key leakage (secure environment binding) with local browser caching for 2G rural networks.")
                ]
                
                for idx, (title, desc) in enumerate(tech_bullets):
                    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                    p.space_after = Pt(6)
                    r1 = p.add_run()
                    r1.text = title
                    r1.font.bold = True
                    r1.font.size = Pt(13)
                    r1.font.color.rgb = RGBColor(15, 23, 42)
                    
                    r2 = p.add_run()
                    r2.text = desc
                    r2.font.bold = False
                    r2.font.size = Pt(12)
                    r2.font.color.rgb = RGBColor(51, 65, 85)

    # ----------------------------------------------------
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ----------------------------------------------------
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_text_frame:
            if 'FEASIBILITY AND VIABILITY' in shape.text:
                shape.text_frame.text = "FEASIBILITY, VIABILITY & OPERATIONAL ROADMAP"
                p = shape.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(16, 185, 129)
            elif 'Your Team Name' in shape.text or 'NeoMe' in shape.text:
                style_oval(shape)
            elif 'Analysis of the feasibility' in shape.text or '1. High Technical' in shape.text:
                shape.top = Inches(1.3)
                shape.left = Inches(0.6)
                shape.width = Inches(11.4)
                shape.height = Inches(5.2)
                tf = shape.text_frame
                tf.clear()
                
                feas_bullets = [
                    ("1. High Technical Feasibility:", " 100% functional working prototype already built and live on Vercel (https://sih-internal-2026-pi.vercel.app/). No dependency on unproven tech."),
                    ("2. Ultra-Low Economic Cost:", " Serverless edge architecture + optimized multimodal inference yields an operational cost of less than ₹0.02 per diagnosis, enabling massive scale across KVKs and Panchayats."),
                    ("3. Potential Risk - Rural Network Drops:", " Mitigated via client-side caching, adaptive image compression, and local benchmark fallback datasets for remote farm fields."),
                    ("4. Potential Risk - Image Glare & Quality Variations:", " Mitigated using camera targeting viewfinders, contrast normalization, and confidence-threshold safety triggers before prescribing chemicals."),
                    ("5. Potential Risk - Dialect & Literacy Barriers:", " Mitigated through native Marwari/Hindi audio synthesis and simplified 2x2 action tile dashboard layout."),
                    ("6. Clinical Safety & Supervision:", " Under clinical domain supervision of AIIMS Jodhpur Asst. Nursing Superintendent, ensuring zero dangerous drug hallucinations.")
                ]
                
                for idx, (title, desc) in enumerate(feas_bullets):
                    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                    p.space_after = Pt(6)
                    r1 = p.add_run()
                    r1.text = title
                    r1.font.bold = True
                    r1.font.size = Pt(13)
                    r1.font.color.rgb = RGBColor(15, 23, 42)
                    
                    r2 = p.add_run()
                    r2.text = desc
                    r2.font.bold = False
                    r2.font.size = Pt(12)
                    r2.font.color.rgb = RGBColor(51, 65, 85)

    # ----------------------------------------------------
    # SLIDE 5: IMPACT AND BENEFITS
    # ----------------------------------------------------
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.has_text_frame:
            if 'IMPACT AND BENEFITS' in shape.text:
                shape.text_frame.text = "IMPACT, SOCIO-ECONOMIC BENEFITS & SUSTAINABILITY"
                p = shape.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(16, 185, 129)
            elif 'Your Team Name' in shape.text or 'NeoMe' in shape.text:
                style_oval(shape)
            elif 'Potential impact on the target audience' in shape.text or '1. Direct Impact' in shape.text:
                shape.top = Inches(1.3)
                shape.left = Inches(0.6)
                shape.width = Inches(11.4)
                shape.height = Inches(5.2)
                tf = shape.text_frame
                tf.clear()
                
                impact_bullets = [
                    ("1. Direct Impact on 140M+ Indian Farmers:", " Delivers instant 24/7 expert veterinary and agricultural triage directly to the farmer's pocket without travel delays or consulting fees."),
                    ("2. Direct Economic Savings (₹4,000–₹8,000/acre/season):", " Prevents catastrophic crop yield losses, stops redundant chemical purchases, and protects high-value livestock capital."),
                    ("3. Livestock Mortality Reduction:", " 60% faster emergency response for contagious outbreaks (Lumpy Skin Disease, FMD) via automated triage and 1962 direct dispatch."),
                    ("4. Environmental & Soil Health Protection:", " Exact spray dosage calculation and organic bio-nutrition recipes reduce chemical pesticide runoff into groundwater by up to 35%."),
                    ("5. Digital Financial Inclusion:", " Built-in Farm & Cattle Ledger allows marginal farmers to record milk yields, fertilizer costs, and track actual seasonal profitability."),
                    ("6. Accessibility & Inclusivity:", " Voice AI in Marwari & Hindi empowers unlettered farmers, women dairy workers, and rural elders to manage farm health independently.")
                ]
                
                for idx, (title, desc) in enumerate(impact_bullets):
                    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                    p.space_after = Pt(6)
                    r1 = p.add_run()
                    r1.text = title
                    r1.font.bold = True
                    r1.font.size = Pt(13)
                    r1.font.color.rgb = RGBColor(15, 23, 42)
                    
                    r2 = p.add_run()
                    r2.text = desc
                    r2.font.bold = False
                    r2.font.size = Pt(12)
                    r2.font.color.rgb = RGBColor(51, 65, 85)

    # ----------------------------------------------------
    # SLIDE 6: RESEARCH AND REFERENCES
    # ----------------------------------------------------
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.has_text_frame:
            if 'RESEARCH' in shape.text:
                shape.text_frame.text = "RESEARCH, REFERENCES & WORKING PROTOTYPE"
                p = shape.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(16, 185, 129)
            elif 'Your Team Name' in shape.text or 'NeoMe' in shape.text:
                style_oval(shape)
            elif 'Details / Links of the reference' in shape.text or '1. ICAR' in shape.text:
                shape.top = Inches(1.3)
                shape.left = Inches(0.6)
                shape.width = Inches(11.4)
                shape.height = Inches(5.2)
                tf = shape.text_frame
                tf.clear()
                
                ref_bullets = [
                    ("1. ICAR & CIBRC Standard Packages of Practices:", " Plant disease chemical formulations, safe dilution ratios (g/L or ml/L), and organic bio-fungicide guidelines (ICAR-CAZRI Jodhpur)."),
                    ("2. AIIMS Jodhpur & VCI Veterinary Protocols:", " Clinical triage frameworks, emergency cattle wound care (KMnO4 antiseptic, turmeric-neem paste), and Department of Animal Husbandry 1962 hotline integration."),
                    ("3. Multimodal Vision Transformer & Neural Diagnostic Research (2026):", " Advanced foliar pathology segmentation and low-latency structured reasoning benchmarks."),
                    ("4. Open-Meteo & IMD Meteorology Datasets:", " Real-time satellite micro-climate telemetry modeling relative humidity (>75%) & heat indices for fungal spore outbreak warnings."),
                    ("5. Live Deployed MVP URL:", " https://sih-internal-2026-pi.vercel.app/ (Fully interactive working application)"),
                    ("6. GitHub Open-Source Codebase:", " https://github.com/maheshiv15/SIH_INTERNAL2026 (Documented source code & build configurations)")
                ]
                
                for idx, (title, desc) in enumerate(ref_bullets):
                    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                    p.space_after = Pt(6)
                    r1 = p.add_run()
                    r1.text = title
                    r1.font.bold = True
                    r1.font.size = Pt(13)
                    r1.font.color.rgb = RGBColor(15, 23, 42)
                    
                    r2 = p.add_run()
                    r2.text = desc
                    r2.font.bold = False
                    r2.font.size = Pt(12)
                    r2.font.color.rgb = RGBColor(51, 65, 85)

    # ----------------------------------------------------
    # REMOVE SLIDE 7 (INSTRUCTIONS SLIDE)
    # ----------------------------------------------------
    if len(prs.slides) > 6:
        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]
        print("Slide 7 (Instructions) successfully removed. Total slides remaining: 6")

    prs.save(output_pptx)
    prs.save(doc_pptx)
    print(f"Presentation saved successfully to: {output_pptx} and {doc_pptx}")

    # Convert to PDF using PowerPoint COM
    try:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        powerpoint.Visible = 1
        deck = powerpoint.Presentations.Open(output_pptx)
        deck.SaveAs(output_pdf, 32) # 32 represents ppSaveAsPDF
        deck.SaveAs(doc_pdf, 32)
        deck.Close()
        powerpoint.Quit()
        pythoncom.CoUninitialize()
        print(f"PDF exported successfully to: {output_pdf} and {doc_pdf}")
    except Exception as e:
        print(f"PowerPoint COM export error: {e}")

if __name__ == '__main__':
    generate_sih_presentation()
