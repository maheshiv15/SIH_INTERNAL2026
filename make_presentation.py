import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    # Set 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Premium Light Aesthetic Palette
    BG_LIGHT = RGBColor(248, 250, 252)       # #f8fafc Clean Soft Canvas
    PANEL_CARD = RGBColor(240, 253, 244)     # #f0fdf4 Light Mint Panel
    BORDER_EMERALD = RGBColor(167, 243, 208) # #a7f3d0 Soft Emerald Border
    PRIMARY_EMERALD = RGBColor(6, 95, 70)    # #065f46 Deep Emerald Headings
    ACCENT_TEAL = RGBColor(13, 148, 136)     # #0d9488 Vibrant Teal Subtitles
    TAG_GREEN = RGBColor(4, 120, 87)         # #047857 Dark Mint Badge
    TEXT_MAIN = RGBColor(30, 41, 59)          # #1e293b Slate Navy Text
    TEXT_MUTED = RGBColor(100, 116, 139)     # #64748b Muted Slate Footer
    WHITE = RGBColor(255, 255, 255)
    
    blank_layout = prs.slide_layouts[6]
    
    # Image Paths
    base_path = r"C:\Users\RAM\.gemini\antigravity-ide\brain\f54736e4-4838-4bd7-8795-a66085400229"
    img_crop = os.path.join(base_path, "home_page_english_1786640419564.png")
    img_livestock = os.path.join(base_path, "livestock_triage_1786640430124.png")
    img_ledger = os.path.join(base_path, "farm_ledger_1786640439309.png")
    img_advisory = os.path.join(base_path, "advisory_alerts_1786640449172.png")
    
    slides_data = [
        {
            "slide_num": 1,
            "tag": "PRODUCTION PLATFORM DEPLOYED",
            "title": "AgriVision AI (Production Platform)",
            "subtitle": "Unified AI Agri-Vision & Veterinary Clinical Diagnostic System",
            "bullets": [
                "Smart India Hackathon 2026 Candidate - IIT Jodhpur",
                "Problem Statement 1: Software Domain",
                "Lead Presenter: Assistant Nursing Superintendent (AIIMS Jodhpur) & MMT Scholar (IIT Jodhpur)",
                "Live Product Status: Fully Deployed Working Product on Vercel & GitHub"
            ],
            "image": img_crop
        },
        {
            "slide_num": 2,
            "tag": "PROBLEM & RURAL IMPACT",
            "title": "The Rural Challenge: Fragmented Care & Loss",
            "subtitle": "Why 120M+ Indian Farmers Suffer Severe Annual Losses",
            "bullets": [
                "Separated Systems: Farmers rely on fragmented crop tools with ZERO digital diagnostic systems for livestock.",
                "Acute Specialist Deficit: Severe shortage of rural agronomists & vets (1 veterinarian per 15,000+ livestock in India).",
                "Catastrophic Shocks: Lumpy Skin Disease & fungal crop blights devastate family incomes within 48 hours."
            ],
            "image": None
        },
        {
            "slide_num": 3,
            "tag": "CORE INNOVATION & ARCHITECTURE",
            "title": "The Solution: Unified Dual-Domain AI Architecture",
            "subtitle": "Single Digital Gateway for Crop Pathology + Livestock Clinical Triage",
            "bullets": [
                "Integrated Doorway: Crop disease detection AND livestock health triage in ONE intuitive interface.",
                "AIIMS Clinical Triage Protocol: Adapting medical emergency triage scoring to animal epidemiology.",
                "Hybrid AI Engine: Real-time Multimodal Neural Vision AI + Offline Rural Field Fallback."
            ],
            "image": img_crop
        },
        {
            "slide_num": 4,
            "tag": "CROP AI PATHOLOGY ENGINE",
            "title": "Live Crop Pathology & ICAR Dosage Calculator",
            "subtitle": "Multi-Image Selection, Live Camera Stream & AI Heatmap Overlay",
            "bullets": [
                "Multi-Angle Photo Inspection: Upload up to 5 leaf photos or capture live stream via Mobile/Webcam.",
                "AI Pathology Heatmap: Visual bounding overlay pinpointing exact leaf lesions and fungal spores.",
                "ICAR Dosage Calculator: Computes exact chemical spray grams/liters & tank refills based on land area."
            ],
            "image": img_crop
        },
        {
            "slide_num": 5,
            "tag": "LIVESTOCK TRIAGE & DIGITAL LEDGER",
            "title": "Livestock Emergency Health Triage & Ledger",
            "subtitle": "AIIMS Nursing Triage Protocol + 1-Click Sync to Digital Herd Ledger",
            "bullets": [
                "Animal Lesion & Symptom Triage: Evaluates cattle, cow, buffalo, goat, & poultry for FMD, Lumpy Skin, & Mastitis.",
                "Emergency Action & Vet SOS: Displays Tier-1 emergency steps & direct call button for 1962 Helpline.",
                "1-Click Digital Ledger Sync: Automatically logs diagnoses and treatments into Farm & Herd Ledger."
            ],
            "image": img_livestock
        },
        {
            "slide_num": 6,
            "tag": "ACCESSIBILITY & VOICE AI",
            "title": "Multilingual Audio Doctor & Conversational AI",
            "subtitle": "Natural Voice Assistance in Hindi, English, and Marwari",
            "bullets": [
                "Conversational Text-To-Speech: Reads out treatment protocols in natural Hindi, English, & Marwari.",
                "Smart Intent Chatbot: Responds intelligently to Hinglish queries ('meri gaay chara nahi kha rahi he', 'chhale/laar').",
                "A4 PDF Official Certificate: Generates downloadable single-page official diagnostic reports with seal."
            ],
            "image": img_advisory
        },
        {
            "slide_num": 7,
            "tag": "COMPETITIVE EDGE",
            "title": "Clinical Rigor Applied to Agriculture",
            "subtitle": "Unique Competitive Advantage of AIIMS & IIT Jodhpur Leadership",
            "bullets": [
                "Medical Triage Precision: Clinical diagnostic workflows applied directly to animal & crop pathology.",
                "Zero False-Negative Goal: Minimizes diagnostic errors in critical infectious outbreaks.",
                "Rural-First UX Design: High-contrast responsive design, large touch targets, and low-bandwidth 2G mode support."
            ],
            "image": img_ledger
        },
        {
            "slide_num": 8,
            "tag": "EPIDEMIC SURVEILLANCE",
            "title": "GIS Outbreak Radar & KVK Early Warning Network",
            "subtitle": "Epidemic Surveillance & Proactive Outbreak Containment",
            "bullets": [
                "District Outbreak Heatmap: Real-time telemetry tracking pest/pathogen cases across Western Rajasthan.",
                "Proactive Radius Alerts: Sends SMS/Voice alerts when Lumpy Skin or Fungal Blight is detected within 15 km.",
                "Prevents Local Spikes: Enables KVKs and animal husbandry departments to halt regional epidemics early."
            ],
            "image": img_advisory
        },
        {
            "slide_num": 9,
            "tag": "PRODUCTION STACK & DEPLOYMENT",
            "title": "Live Production Stack & Deployment Model",
            "subtitle": "Fully Built, Tested & Deployed Today",
            "bullets": [
                "Production Stack: React 18, Vite, Vanilla CSS, Lucide Icons, Canvas API, WebRTC MediaDevices.",
                "AI Services: Multimodal Neural Vision Model API + Local Edge Diagnostics Engine.",
                "Live Deployment: Hosted live on Vercel with automatic GitHub CI/CD integration."
            ],
            "image": img_ledger
        },
        {
            "slide_num": 10,
            "tag": "CONCLUSION & Q&A",
            "title": "Summary & Call to Action",
            "subtitle": "Ready for SIH 2026 Internal Hackathon Evaluation",
            "bullets": [
                "100% Mandate Fulfilled: Crop AI, Livestock Triage, Farm Records, Advisory Services, & Outbreak Radar in one app.",
                "Fully Functional Live Product: Tested with real crop leaves, cattle lesion photos, and multi-lingual voice audio.",
                "Thank You! Team AgriVision AI (IIT Jodhpur) is Ready for Evaluator Q&A & Live Demonstration."
            ],
            "image": img_crop
        }
    ]

    for data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Background shape - Soft Pristine Canvas
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_LIGHT
        bg.line.fill.background()
        
        # Top Accent Decorative Bar
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.06))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = TAG_GREEN
        top_bar.line.fill.background()

        # Top Tag Badge
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.35))
        tf_tag = tag_box.text_frame
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = f"SLIDE {data['slide_num']}  |  {data['tag']}"
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = TAG_GREEN
        
        # Main Title (Deep Emerald)
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.7), Inches(0.75))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = data['title']
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = PRIMARY_EMERALD
        
        # Subtitle (Teal Blue)
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.45))
        tf_sub = sub_box.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = data['subtitle']
        p_sub.font.size = Pt(15)
        p_sub.font.bold = True
        p_sub.font.color.rgb = ACCENT_TEAL
        
        has_image = data['image'] and os.path.exists(data['image'])
        text_width = Inches(5.8) if has_image else Inches(11.733)
        
        # Bullets Panel Card (Soft Mint Background with Emerald Border)
        bullet_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), text_width, Inches(4.6))
        bullet_panel.fill.solid()
        bullet_panel.fill.fore_color.rgb = PANEL_CARD
        bullet_panel.line.color.rgb = BORDER_EMERALD
        bullet_panel.line.width = Pt(1.5)
        
        # Bullet Text
        tb_bullets = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), text_width - Inches(0.4), Inches(4.2))
        tf_b = tb_bullets.text_frame
        tf_b.word_wrap = True
        
        for i, bullet in enumerate(data['bullets']):
            p = tf_b.paragraphs[0] if i == 0 else tf_b.add_paragraph()
            p.text = f"✔  {bullet}"
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_MAIN
            p.space_after = Pt(14)
            p.line_spacing = 1.25
            
        # Add App Screenshot if available (Clean Framed Shadow Look)
        if has_image:
            img_left = Inches(6.933)
            img_top = Inches(2.2)
            img_width = Inches(5.6)
            img_height = Inches(4.6)
            
            # Screenshot Border Frame Card
            img_border = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, img_left - Inches(0.04), img_top - Inches(0.04), img_width + Inches(0.08), img_height + Inches(0.08))
            img_border.fill.solid()
            img_border.fill.fore_color.rgb = WHITE
            img_border.line.color.rgb = BORDER_EMERALD
            img_border.line.width = Pt(2)
            
            slide.shapes.add_picture(data['image'], img_left, img_top, width=img_width, height=img_height)
            
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.95), Inches(11.7), Inches(0.3))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "AgriVision AI Platform  •  SIH 2026 Internal Hackathon  •  IIT Jodhpur & AIIMS Jodhpur"
        p_foot.font.size = Pt(10)
        p_foot.font.color.rgb = TEXT_MUTED

    output_path = r"d:\NEW CODING PROJECTS\SIH_INTERNAL\AgriVision_AI_SIH2026_Presentation.pptx"
    prs.save(output_path)
    print(f"SUCCESS: Created Light Theme PowerPoint presentation at {output_path}")

if __name__ == "__main__":
    build_presentation()
